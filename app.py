import os
import io
import json
import tempfile
import requests
import re

from flask import Flask, render_template, request, send_file, redirect, url_for, flash
from werkzeug.utils import secure_filename
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

app = Flask(__name__, static_folder='static')
app.config['MAX_CONTENT_LENGTH'] = 8 * 1024 * 1024  # 8MB limit
app.secret_key = 'your-secret-key'  # For flash messages

UPLOAD_FOLDER = tempfile.gettempdir()

def heuristic_slide_json(text):
    paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
    slides = []
    if not paragraphs:
        return {'slides': []}
    first_para = paragraphs[0]
    lines = first_para.split('\n')
    if len(lines) > 1:
        title, subtitle = lines[0], lines[1]
    else:
        period = first_para.find('.')
        if period != -1:
            title, subtitle = first_para[:period], first_para[period+1:]
        else:
            title, subtitle = first_para, ''
    slides.append({'type': 'title_slide', 'title': title.strip(), 'subtitle': subtitle.strip()})
    for para in paragraphs[1:]:
        lines = para.split('\n')
        slide_title = lines[0].strip()
        if len(lines) > 1:
            content = [line.strip() for line in lines[1:] if line.strip()]
        else:
            content = [para.strip()]
        slides.append({'type': 'content_slide', 'title': slide_title, 'content': content})
    return {'slides': slides}

def call_llm_api(api_key, provider, text, guidance):
    endpoint = ''
    headers = {}
    payload = {}
    if provider.lower() == 'openai':
        endpoint = 'https://api.openai.com/v1/chat/completions'
        headers = {
            'Authorization': f'Bearer {api_key}',
            'Content-Type': 'application/json'
        }
        sys_prompt = (
            "You are an expert presentation creator. Given user text and guidance, "
            "break the content into a practical number of slides for a PowerPoint deck. "
            "Return ONLY a valid JSON object according to this schema:\n"
            "{\"slides\": [{\"type\": \"title_slide\", \"title\": \"...\", \"subtitle\": \"...\"},"
            "{\"type\": \"content_slide\", \"title\": \"...\", \"content\": [\"...\"]},"
            "{\"type\": \"image_content_slide\", \"title\": \"...\", \"content\": [\"...\"], \"image_suggestion\": \"...\"}]}"
        )
        payload = {
            'model': 'gpt-3.5-turbo',
            'messages': [
                {'role': 'system', 'content': sys_prompt},
                {'role': 'user', 'content': f'Text:\n{text}\nGuidance: {guidance}'}
            ],
            'temperature': 0.3
        }
    else:
        return {'slides': []}
    try:
        response = requests.post(endpoint, headers=headers, data=json.dumps(payload), timeout=35)
        res_data = response.json()
        for choice in res_data.get('choices', []):
            content = choice['message']['content'].strip()
            match = re.search(r'\{[\s\S]*\}', content)
            if match:
                json_str = match.group(0)
            else:
                json_str = re.sub(r'``````', r'\1', content, flags=re.DOTALL).strip()
            try:
                parsed = json.loads(json_str)
                if 'slides' in parsed:
                    return parsed
            except Exception:
                continue
    except Exception:
        pass
    return {'slides': []}

def extract_template_images(prs):
    images = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                ext = image.ext if image.ext else "png"
                images.append({'data': image.blob, 'ext': ext})
    return images

def find_layout_index(prs, target_names):
    for layout in prs.slide_layouts:
        if layout.name in target_names:
            return prs.slide_layouts.index(layout)
    return 0  # fallback

def clear_text_frame(tf):
    p_count = len(tf.paragraphs)
    for idx in reversed(range(p_count)):
        p = tf.paragraphs[idx]
        tf._element.remove(p._element)

def add_text_to_placeholder(ph, lines, font_size=18):
    if not ph.has_text_frame:
        return
    tf = ph.text_frame
    clear_text_frame(tf)
    if isinstance(lines, str):
        lines = [lines]
    for line in lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

def add_textbox(slide, left, top, width, height, lines, font_size=18):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    clear_text_frame(tf)
    if isinstance(lines, str):
        lines = [lines]
    for line in lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    return txBox

def add_slide_from_json(prs, layout_idx, slide_json, template_images):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    slide_type = slide_json.get('type', 'content_slide')

    # TITLE
    title_text = slide_json.get('title', '')
    title_ph = None
    for ph in slide.placeholders:
        if ph.placeholder_format.type == PP_PLACEHOLDER.TITLE:
            title_ph = ph
            break
    if title_ph:
        add_text_to_placeholder(title_ph, title_text, font_size=24)
    else:
        add_textbox(slide, Inches(0.5), Inches(0.1), Inches(9), Inches(1), title_text, font_size=24)

    if slide_type == 'title_slide':
        subtitle_text = slide_json.get('subtitle', '')
        subtitle_ph = None
        for ph in slide.placeholders:
            if ph.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE:
                subtitle_ph = ph
                break
        if subtitle_ph:
            add_text_to_placeholder(subtitle_ph, subtitle_text, font_size=18)
        else:
            add_textbox(slide, Inches(0.5), Inches(1.1), Inches(9), Inches(1), subtitle_text, font_size=18)
    else:
        content_lines = slide_json.get('content', [])
        if not isinstance(content_lines, list):
            content_lines = [content_lines]
        body_ph = None
        for ph in slide.placeholders:
            if ph.placeholder_format.type == PP_PLACEHOLDER.BODY:
                body_ph = ph
                break
        if body_ph:
            add_text_to_placeholder(body_ph, content_lines, font_size=18)
        else:
            add_textbox(slide, Inches(0.5), Inches(1), Inches(9), Inches(4), content_lines, font_size=18)

        if slide_type == 'image_content_slide' and template_images:
            img = template_images[0]
            tmp_img_path = os.path.join(tempfile.gettempdir(), f'tmp_img.{img["ext"]}')
            with open(tmp_img_path, 'wb') as f:
                f.write(img['data'])
            slide.shapes.add_picture(tmp_img_path, Inches(6.5), Inches(2), width=Inches(3))
            os.remove(tmp_img_path)

    return slide

@app.route('/', methods=['GET'])
def index():
    return render_template('index.html')

@app.route('/generate', methods=['POST'])
def generate():
    text = request.form.get('source_text', '').strip()
    guidance = request.form.get('guidance', '').strip()
    api_key = request.form.get('api_key', '').strip()
    provider = request.form.get('api_provider', 'openai').strip()
    file = request.files.get('template_file')
    if not text or not file:
        flash('Please provide both text and a template file!', 'error')
        return redirect(url_for('index'))
    filename = secure_filename(file.filename)
    template_path = os.path.join(UPLOAD_FOLDER, filename)
    file.save(template_path)

    template_prs = Presentation(template_path)
    images = extract_template_images(template_prs)
    prs = Presentation(template_path)

    title_slide_idx = find_layout_index(prs, ['Title Slide', 'Title'])
    content_slide_idx = find_layout_index(prs, ['Title and Content', 'Content'])
    image_slide_idx = find_layout_index(prs, ['Picture with Caption', 'Title and Content', 'Content'])

    if api_key:
        slide_json = call_llm_api(api_key, provider, text, guidance)
    else:
        slide_json = heuristic_slide_json(text)

    for slide in slide_json.get('slides', []):
        if slide.get('type') == 'title_slide':
            layout_idx = title_slide_idx
        elif slide.get('type') == 'content_slide':
            layout_idx = content_slide_idx
        elif slide.get('type') == 'image_content_slide':
            layout_idx = image_slide_idx
        else:
            layout_idx = content_slide_idx
        add_slide_from_json(prs, layout_idx, slide, images)

    output_stream = io.BytesIO()
    prs.save(output_stream)
    output_stream.seek(0)
    os.remove(template_path)
    return send_file(
        output_stream,
        as_attachment=True,
        download_name="generated_presentation.pptx",
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

if __name__ == '__main__':
    app.run(debug=True)
